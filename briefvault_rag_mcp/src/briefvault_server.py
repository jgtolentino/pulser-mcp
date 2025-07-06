#!/usr/bin/env python3
"""
BriefVault RAG MCP Server - Complex Document Processing
Advanced RAG system for marketing briefs, strategy docs, and creative assets
"""

import os
import json
import uuid
from pathlib import Path
from typing import Dict, List, Any, Optional
from datetime import datetime
import logging
from fastapi import FastAPI, HTTPException, Depends, UploadFile, File
from pydantic import BaseModel
import uvicorn
import hashlib
from qdrant_client import QdrantClient
from qdrant_client.models import Distance, VectorParams, PointStruct
import fitz  # PyMuPDF
from PIL import Image
import pytesseract
import docx
from pptx import Presentation
import pandas as pd
from sentence_transformers import SentenceTransformer
import re
from transformers import BlipProcessor, BlipForConditionalGeneration
import torch

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# FastAPI app

from fastapi import Depends, HTTPException, status
from fastapi.security import HTTPBearer, HTTPAuthorizationCredentials
from jose import JWTError, jwt
from datetime import datetime, timedelta
from typing import Optional
import os

# JWT Configuration
SECRET_KEY = os.getenv("PULSER_JWT_SECRET", "change-this-in-production")
ALGORITHM = "HS256"
ACCESS_TOKEN_EXPIRE_MINUTES = 30

security = HTTPBearer()

def create_access_token(data: dict, expires_delta: Optional[timedelta] = None):
    """Create JWT access token"""
    to_encode = data.copy()
    if expires_delta:
        expire = datetime.utcnow() + expires_delta
    else:
        expire = datetime.utcnow() + timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES)
    to_encode.update({"exp": expire})
    encoded_jwt = jwt.encode(to_encode, SECRET_KEY, algorithm=ALGORITHM)
    return encoded_jwt

async def verify_token(credentials: HTTPAuthorizationCredentials = Depends(security)):
    """Verify JWT token"""
    token = credentials.credentials
    try:
        payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
        username: str = payload.get("sub")
        if username is None:
            raise HTTPException(
                status_code=status.HTTP_401_UNAUTHORIZED,
                detail="Invalid authentication credentials",
                headers={"WWW-Authenticate": "Bearer"},
            )
        return username
    except JWTError:
        raise HTTPException(
            status_code=status.HTTP_401_UNAUTHORIZED,
            detail="Invalid authentication credentials",
            headers={"WWW-Authenticate": "Bearer"},
        )

# Public endpoints that don't require authentication
PUBLIC_ENDPOINTS = ["/", "/health", "/auth/token"]

app = FastAPI(title="BriefVault RAG MCP Server")

from fastapi import APIRouter

# Create versioned router
api_v1 = APIRouter(prefix="/api/v1")

# Data directories
DATA_DIR = Path(__file__).parent.parent / "data"
DOCS_DIR = DATA_DIR / "documents"
PROCESSED_DIR = DATA_DIR / "processed"
LOGS_DIR = Path(__file__).parent.parent / "logs"

for dir_path in [DATA_DIR, DOCS_DIR, PROCESSED_DIR, LOGS_DIR]:
    dir_path.mkdir(exist_ok=True)

# Initialize models and clients
text_encoder = SentenceTransformer('all-MiniLM-L6-v2')
qdrant_client = QdrantClient(":memory:")  # Use in-memory for demo

# Initialize image captioning model
try:
    blip_processor = BlipProcessor.from_pretrained("Salesforce/blip-image-captioning-base")
    blip_model = BlipForConditionalGeneration.from_pretrained("Salesforce/blip-image-captioning-base")
    logger.info("Image captioning model loaded")
except Exception as e:
    logger.warning(f"Could not load image model: {e}")
    blip_processor = None
    blip_model = None

# Pydantic models
class DocumentIngestionRequest(BaseModel):
    file_path: str
    document_type: str = "unknown"  # creative_brief, strategy_doc, brand_guidelines, etc.
    metadata: Optional[Dict[str, Any]] = {}

class SearchRequest(BaseModel):
    query: str
    filters: Optional[Dict[str, Any]] = {}
    limit: int = 10
    include_snippets: bool = True

class AnalysisRequest(BaseModel):
    document_id: str
    analysis_type: List[str] = ["objectives", "audience", "key_messages"]

class ComplianceCheckRequest(BaseModel):
    asset_path: str
    brand_guidelines_id: str

# Document processing utilities
class DocumentProcessor:
    @staticmethod
    def extract_text_from_pdf(file_path: str) -> Dict[str, Any]:
        """Extract text and metadata from PDF"""
        try:
            doc = fitz.open(file_path)
            content = {
                "text": "",
                "pages": [],
                "images": [],
                "metadata": doc.metadata
            }
            
            for page_num in range(len(doc)):
                page = doc[page_num]
                page_text = page.get_text()
                content["text"] += page_text
                content["pages"].append({
                    "page_number": page_num + 1,
                    "text": page_text,
                    "bbox": page.rect
                })
                
                # Extract images
                image_list = page.get_images()
                for img_index, img in enumerate(image_list):
                    try:
                        xref = img[0]
                        pix = fitz.Pixmap(doc, xref)
                        if pix.n < 5:  # GRAY or RGB
                            img_data = pix.tobytes("png")
                            content["images"].append({
                                "page": page_num + 1,
                                "index": img_index,
                                "data": img_data,
                                "bbox": img[1] if len(img) > 1 else None
                            })
                        pix = None
                    except Exception as e:
                        logger.warning(f"Could not extract image {img_index} from page {page_num}: {e}")
            
            doc.close()
            return content
            
        except Exception as e:
            logger.error(f"Error processing PDF {file_path}: {e}")
            return {"text": "", "pages": [], "images": [], "metadata": {}}
    
    @staticmethod
    def extract_text_from_docx(file_path: str) -> Dict[str, Any]:
        """Extract text and metadata from Word document"""
        try:
            doc = docx.Document(file_path)
            content = {
                "text": "",
                "paragraphs": [],
                "tables": [],
                "metadata": {}
            }
            
            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    content["text"] += para.text + "\n"
                    content["paragraphs"].append({
                        "text": para.text,
                        "style": para.style.name if para.style else None
                    })
            
            # Extract tables
            for table_idx, table in enumerate(doc.tables):
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                content["tables"].append({
                    "index": table_idx,
                    "data": table_data
                })
            
            # Extract document properties
            props = doc.core_properties
            content["metadata"] = {
                "title": props.title,
                "author": props.author,
                "created": props.created.isoformat() if props.created else None,
                "modified": props.modified.isoformat() if props.modified else None
            }
            
            return content
            
        except Exception as e:
            logger.error(f"Error processing DOCX {file_path}: {e}")
            return {"text": "", "paragraphs": [], "tables": [], "metadata": {}}
    
    @staticmethod
    def extract_text_from_pptx(file_path: str) -> Dict[str, Any]:
        """Extract text and metadata from PowerPoint presentation"""
        try:
            prs = Presentation(file_path)
            content = {
                "text": "",
                "slides": [],
                "metadata": {}
            }
            
            for slide_idx, slide in enumerate(prs.slides):
                slide_content = {
                    "slide_number": slide_idx + 1,
                    "title": "",
                    "text": "",
                    "notes": ""
                }
                
                # Extract slide content
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.strip()
                        slide_content["text"] += text + "\n"
                        content["text"] += text + "\n"
                        
                        # Try to identify title
                        if shape.shape_type == 1 and not slide_content["title"]:  # Title placeholder
                            slide_content["title"] = text
                
                # Extract speaker notes
                if slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text
                    if notes_text.strip():
                        slide_content["notes"] = notes_text.strip()
                        content["text"] += notes_text + "\n"
                
                content["slides"].append(slide_content)
            
            return content
            
        except Exception as e:
            logger.error(f"Error processing PPTX {file_path}: {e}")
            return {"text": "", "slides": [], "metadata": {}}

class BriefAnalyzer:
    @staticmethod
    def extract_objectives(text: str) -> List[str]:
        """Extract campaign objectives from brief text"""
        objectives = []
        
        # Common objective patterns
        patterns = [
            r"objective[s]?[:\-\s]+(.*?)(?:\n\n|\n[A-Z]|\n\d+\.|\Z)",
            r"goal[s]?[:\-\s]+(.*?)(?:\n\n|\n[A-Z]|\n\d+\.|\Z)",
            r"aim[s]?[:\-\s]+(.*?)(?:\n\n|\n[A-Z]|\n\d+\.|\Z)",
            r"purpose[:\-\s]+(.*?)(?:\n\n|\n[A-Z]|\n\d+\.|\Z)"
        ]
        
        for pattern in patterns:
            matches = re.finditer(pattern, text, re.IGNORECASE | re.DOTALL)
            for match in matches:
                obj = match.group(1).strip()
                if len(obj) > 10 and len(obj) < 500:  # Reasonable length
                    objectives.append(obj)
        
        return objectives[:5]  # Return top 5
    
    @staticmethod
    def extract_target_audience(text: str) -> Dict[str, Any]:
        """Extract target audience information"""
        audience = {
            "demographics": [],
            "psychographics": [],
            "behaviors": []
        }
        
        # Demographic patterns
        demo_patterns = [
            r"target audience[:\-\s]+(.*?)(?:\n\n|\n[A-Z]|\Z)",
            r"demographics?[:\-\s]+(.*?)(?:\n\n|\n[A-Z]|\Z)",
            r"age[:\-\s]+(\d+[^\n]*)",
            r"income[:\-\s]+([^\n]*)",
            r"gender[:\-\s]+([^\n]*)"
        ]
        
        for pattern in demo_patterns:
            matches = re.finditer(pattern, text, re.IGNORECASE | re.DOTALL)
            for match in matches:
                demo = match.group(1).strip()
                if len(demo) > 5:
                    audience["demographics"].append(demo)
        
        return audience
    
    @staticmethod
    def extract_key_messages(text: str) -> List[str]:
        """Extract key messages and positioning"""
        messages = []
        
        patterns = [
            r"key message[s]?[:\-\s]+(.*?)(?:\n\n|\n[A-Z]|\n\d+\.|\Z)",
            r"main message[s]?[:\-\s]+(.*?)(?:\n\n|\n[A-Z]|\n\d+\.|\Z)",
            r"positioning[:\-\s]+(.*?)(?:\n\n|\n[A-Z]|\n\d+\.|\Z)",
            r"value proposition[:\-\s]+(.*?)(?:\n\n|\n[A-Z]|\n\d+\.|\Z)"
        ]
        
        for pattern in patterns:
            matches = re.finditer(pattern, text, re.IGNORECASE | re.DOTALL)
            for match in matches:
                msg = match.group(1).strip()
                if len(msg) > 10 and len(msg) < 300:
                    messages.append(msg)
        
        return messages[:3]  # Return top 3

# Initialize Qdrant collection
def setup_qdrant():
    """Initialize Qdrant collection for document storage"""
    try:
        collections = qdrant_client.get_collections().collections
        collection_names = [col.name for col in collections]
        
        if "briefvault" not in collection_names:
            qdrant_client.create_collection(
                collection_name="briefvault",
                vectors_config=VectorParams(size=384, distance=Distance.COSINE)
            )
            logger.info("Created BriefVault collection in Qdrant")
        else:
            logger.info("BriefVault collection already exists")
            
    except Exception as e:
        logger.error(f"Error setting up Qdrant: {e}")

# MCP Tools
class MCPTools:
    @staticmethod
    async def ingest_document(request: DocumentIngestionRequest) -> Dict[str, Any]:
        """Ingest and process a complex document"""
        try:
            file_path = Path(request.file_path)
            if not file_path.exists():
                return {"success": False, "error": f"File not found: {file_path}"}
            
            # Determine file type and process accordingly
            extension = file_path.suffix.lower()
            
            if extension == '.pdf':
                content = DocumentProcessor.extract_text_from_pdf(str(file_path))
            elif extension == '.docx':
                content = DocumentProcessor.extract_text_from_docx(str(file_path))
            elif extension == '.pptx':
                content = DocumentProcessor.extract_text_from_pptx(str(file_path))
            else:
                return {"success": False, "error": f"Unsupported file type: {extension}"}
            
            if not content.get("text"):
                return {"success": False, "error": "No text content extracted"}
            
            # Generate document ID
            doc_id = hashlib.md5(f"{file_path.name}{datetime.now().isoformat()}".encode()).hexdigest()
            
            # Create embeddings
            text_chunks = MCPTools._chunk_text(content["text"])
            embeddings = []
            
            for i, chunk in enumerate(text_chunks):
                embedding = text_encoder.encode(chunk).tolist()
                
                # Create metadata
                metadata = {
                    "document_id": doc_id,
                    "chunk_id": i,
                    "file_name": file_path.name,
                    "file_path": str(file_path),
                    "document_type": request.document_type,
                    "processed_at": datetime.now().isoformat(),
                    "chunk_text": chunk,
                    **request.metadata,
                    **content.get("metadata", {})
                }
                
                embeddings.append({
                    "id": f"{doc_id}_{i}",
                    "vector": embedding,
                    "metadata": metadata
                })
            
            # Store in Qdrant
            points = [
                PointStruct(
                    id=emb["id"],
                    vector=emb["vector"],
                    payload=emb["metadata"]
                )
                for emb in embeddings
            ]
            
            qdrant_client.upsert(
                collection_name="briefvault",
                points=points
            )
            
            # Save processed document
            processed_doc = {
                "document_id": doc_id,
                "original_path": str(file_path),
                "content": content,
                "metadata": request.metadata,
                "processed_at": datetime.now().isoformat(),
                "chunk_count": len(text_chunks)
            }
            
            processed_file = PROCESSED_DIR / f"{doc_id}.json"
            with open(processed_file, 'w') as f:
                json.dump(processed_doc, f, indent=2, default=str)
            
            logger.info(f"Processed document {file_path.name} into {len(text_chunks)} chunks")
            
            return {
                "success": True,
                "document_id": doc_id,
                "chunks_created": len(text_chunks),
                "file_type": extension,
                "processed_file": str(processed_file)
            }
            
        except Exception as e:
            logger.error(f"Error ingesting document: {e}")
            return {"success": False, "error": str(e)}
    
    @staticmethod
    def _chunk_text(text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
        """Split text into overlapping chunks"""
        chunks = []
        start = 0
        
        while start < len(text):
            end = start + chunk_size
            chunk = text[start:end]
            
            # Try to break at sentence boundary
            if end < len(text):
                last_period = chunk.rfind('.')
                last_newline = chunk.rfind('\n')
                break_point = max(last_period, last_newline)
                
                if break_point > start + chunk_size // 2:
                    chunk = chunk[:break_point + 1]
                    end = start + len(chunk)
            
            chunks.append(chunk.strip())
            start = end - overlap
            
            if start >= len(text):
                break
                
        return [chunk for chunk in chunks if len(chunk.strip()) > 50]
    
    @staticmethod
    async def search_briefs(request: SearchRequest) -> Dict[str, Any]:
        """Search through processed briefs using semantic similarity"""
        try:
            # Generate query embedding
            query_embedding = text_encoder.encode(request.query).tolist()
            
            # Build search filters
            search_filter = None
            if request.filters:
                conditions = []
                for key, value in request.filters.items():
                    if isinstance(value, list):
                        conditions.extend([{"key": key, "match": {"value": v}} for v in value])
                    else:
                        conditions.append({"key": key, "match": {"value": value}})
                
                if conditions:
                    search_filter = {"must": conditions}
            
            # Search in Qdrant
            search_results = qdrant_client.search(
                collection_name="briefvault",
                query_vector=query_embedding,
                query_filter=search_filter,
                limit=request.limit,
                with_payload=True
            )
            
            results = []
            for result in search_results:
                result_data = {
                    "score": result.score,
                    "document_id": result.payload.get("document_id"),
                    "file_name": result.payload.get("file_name"),
                    "document_type": result.payload.get("document_type"),
                    "metadata": {k: v for k, v in result.payload.items() 
                              if k not in ["chunk_text", "document_id", "chunk_id"]}
                }
                
                if request.include_snippets:
                    result_data["snippet"] = result.payload.get("chunk_text", "")[:300] + "..."
                
                results.append(result_data)
            
            return {
                "success": True,
                "query": request.query,
                "results": results,
                "total_results": len(results)
            }
            
        except Exception as e:
            logger.error(f"Error searching briefs: {e}")
            return {"success": False, "error": str(e)}
    
    @staticmethod
    async def analyze_brief(request: AnalysisRequest) -> Dict[str, Any]:
        """Analyze a brief for key components"""
        try:
            # Load processed document
            processed_file = PROCESSED_DIR / f"{request.document_id}.json"
            if not processed_file.exists():
                return {"success": False, "error": f"Document {request.document_id} not found"}
            
            with open(processed_file, 'r') as f:
                doc_data = json.load(f)
            
            text = doc_data["content"]["text"]
            analysis = {}
            
            # Perform requested analysis
            if "objectives" in request.analysis_type:
                analysis["objectives"] = BriefAnalyzer.extract_objectives(text)
            
            if "audience" in request.analysis_type:
                analysis["target_audience"] = BriefAnalyzer.extract_target_audience(text)
            
            if "key_messages" in request.analysis_type:
                analysis["key_messages"] = BriefAnalyzer.extract_key_messages(text)
            
            if "creative_requirements" in request.analysis_type:
                # Extract creative requirements (simplified)
                requirements = []
                creative_patterns = [
                    r"creative requirements?[:\-\s]+(.*?)(?:\n\n|\n[A-Z]|\Z)",
                    r"deliverables?[:\-\s]+(.*?)(?:\n\n|\n[A-Z]|\Z)",
                    r"assets? needed[:\-\s]+(.*?)(?:\n\n|\n[A-Z]|\Z)"
                ]
                
                for pattern in creative_patterns:
                    matches = re.finditer(pattern, text, re.IGNORECASE | re.DOTALL)
                    for match in matches:
                        req = match.group(1).strip()
                        if len(req) > 10:
                            requirements.append(req)
                
                analysis["creative_requirements"] = requirements[:5]
            
            return {
                "success": True,
                "document_id": request.document_id,
                "analysis": analysis,
                "analyzed_at": datetime.now().isoformat()
            }
            
        except Exception as e:
            logger.error(f"Error analyzing brief: {e}")
            return {"success": False, "error": str(e)}

# API Endpoints
@app.get("/")
async def root():
    return {
        "service": "BriefVault RAG MCP Server",
        "version": "1.0.0",
        "status": "running",
        "features": [
            "Complex document processing (PDF, DOCX, PPTX)",
            "Multi-modal content extraction",
            "Semantic search with embeddings",
            "Brief intelligence analysis",
            "Creative compliance checking"
        ]
    }


@app.post("/auth/token")
async def login(username: str, password: str):
    """Authenticate and get access token"""
    # In production, verify against secure user store
    if username == os.getenv("MCP_ADMIN_USER", "admin") and password == os.getenv("MCP_ADMIN_PASS", "change-this"):
        access_token_expires = timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES)
        access_token = create_access_token(
            data={"sub": username}, expires_delta=access_token_expires
        )
        return {"access_token": access_token, "token_type": "bearer"}
    raise HTTPException(
        status_code=status.HTTP_401_UNAUTHORIZED,
        detail="Incorrect username or password",
        headers={"WWW-Authenticate": "Bearer"},
    )

@app.get("/health")
async def health():
    try:
        collections = qdrant_client.get_collections()
        collection_count = len(collections.collections)
        
        # Count documents
        briefvault_info = qdrant_client.get_collection("briefvault")
        doc_count = briefvault_info.points_count
        
        return {
            "status": "healthy",
            "timestamp": datetime.now().isoformat(),
            "collections": collection_count,
            "documents_indexed": doc_count,
            "models_loaded": {
                "text_encoder": text_encoder is not None,
                "image_model": blip_model is not None
            }
        }
    except Exception as e:
        return {
            "status": "unhealthy",
            "error": str(e),
            "timestamp": datetime.now().isoformat()
        }

@api_v1.post("/mcp/tools/ingest_document")
async def ingest_document(request: DocumentIngestionRequest, current_user: str = Depends(verify_token)):
    """Ingest and process a complex document"""
    return await MCPTools.ingest_document(request)

@api_v1.post("/mcp/tools/search_briefs")
async def search_briefs(request: SearchRequest, current_user: str = Depends(verify_token)):
    """Search through processed briefs"""
    return await MCPTools.search_briefs(request)

@api_v1.post("/mcp/tools/analyze_brief")
async def analyze_brief(request: AnalysisRequest, current_user: str = Depends(verify_token)):
    """Analyze a brief for key components"""
    return await MCPTools.analyze_brief(request)

@api_v1.post("/mcp/tools/upload_document")
async def upload_document(
    file: UploadFile = File(...),
    document_type: str = "unknown",
    metadata: str = "{}"
, current_user: str = Depends(verify_token)):
    """Upload and process a document"""
    try:
        # Save uploaded file
        file_path = DOCS_DIR / file.filename
        with open(file_path, "wb") as buffer:
            content = await file.read()
            buffer.write(content)
        
        # Process the document
        request = DocumentIngestionRequest(
            file_path=str(file_path),
            document_type=document_type,
            metadata=json.loads(metadata) if metadata != "{}" else {}
        )
        
        result = await MCPTools.ingest_document(request)
        
        return result
        
    except Exception as e:
        logger.error(f"Error uploading document: {e}")
        return {"success": False, "error": str(e)}


# Include API v1 router
app.include_router(api_v1)

if __name__ == "__main__":
    logger.info("Starting BriefVault RAG MCP Server...")
    
    # Setup Qdrant
    setup_qdrant()
    
    logger.info(f"Data directory: {DATA_DIR}")
    logger.info(f"Documents directory: {DOCS_DIR}")
    logger.info("Server running at http://localhost:8006")
    
    uvicorn.run(app, host="0.0.0.0", port=8006)